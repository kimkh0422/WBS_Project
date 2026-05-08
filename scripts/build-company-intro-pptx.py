"""
지엠티 프로젝트 매니저 회사 도입·배포용 사용자 안내 PPT 생성.

특징:
  - 16:9 와이드, 한국어, 10슬라이드 이내
  - 배포용(혼자 보고 이해 가능) — 슬라이드별 본문이 자족적
  - 이미지 자리는 회색 박스 + 캡처 안내 문구
  - 회사 도입 관점: 가치 → 사용 흐름 → 협업

실행:
  python scripts/build-company-intro-pptx.py
출력:
  docs/manual/지엠티스마트시트_도입안내.pptx
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# ────────────────────────────────────────────────────────────────────────────
# 디자인 토큰
# ────────────────────────────────────────────────────────────────────────────
COLOR_INK = RGBColor(0x1F, 0x29, 0x37)        # 본문 (어두운 슬레이트)
COLOR_INK_MUTED = RGBColor(0x6B, 0x72, 0x80)  # 보조
COLOR_ACCENT = RGBColor(0x3B, 0x82, 0xF6)     # 강조 (블루)
COLOR_ACCENT_DARK = RGBColor(0x1E, 0x40, 0xAF)
COLOR_HIGHLIGHT = RGBColor(0xEF, 0x44, 0x44)  # 빨강 (캡처 안내)
COLOR_BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
COLOR_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
COLOR_PLACEHOLDER = RGBColor(0xF3, 0xF4, 0xF6)
COLOR_PLACEHOLDER_TEXT = RGBColor(0x9C, 0xA3, 0xAF)
COLOR_TIP_BG = RGBColor(0xFF, 0xF8, 0xE1)
COLOR_TIP_BORDER = RGBColor(0xFC, 0xD3, 0x4D)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_KO = "맑은 고딕"


# ────────────────────────────────────────────────────────────────────────────
# 헬퍼
# ────────────────────────────────────────────────────────────────────────────

def _set_run(run, *, font_name: str = FONT_KO, size: int = 14, bold: bool = False,
             color: RGBColor | None = None, italic: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, *, left, top, width, height, text: str,
                size: int = 14, bold: bool = False,
                color: RGBColor = COLOR_INK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set_run(r, size=size, bold=bold, color=color)
    return tb


def add_filled_rect(slide, *, left, top, width, height,
                    fill: RGBColor, line: RGBColor | None = None,
                    line_width: float = 0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_image_placeholder(slide, *, left, top, width, height,
                          label: str, hint: str):
    """이미지 자리: 회색 박스 + 라벨 + 캡처 안내 텍스트."""
    box = add_filled_rect(
        slide, left=left, top=top, width=width, height=height,
        fill=COLOR_PLACEHOLDER, line=COLOR_BORDER, line_width=1.0,
    )
    # 라벨 (📷 아이콘 + 번호)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(150000)
    tf.margin_right = Emu(150000)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    _set_run(r1, size=14, bold=True, color=COLOR_PLACEHOLDER_TEXT)
    # 안내
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = hint
    _set_run(r2, size=11, color=COLOR_PLACEHOLDER_TEXT)
    return box


def add_slide_header(slide, *, page_no: int, total: int, title: str, subtitle: str | None = None,
                     slide_w: int = 0, slide_h: int = 0):
    """모든 컨텐츠 슬라이드 공통 헤더: 좌측 제목, 우측 페이지 번호."""
    # 좌측 강조 막대
    add_filled_rect(slide, left=Inches(0.5), top=Inches(0.45),
                    width=Inches(0.06), height=Inches(0.5),
                    fill=COLOR_ACCENT, line=None)
    # 제목
    add_textbox(slide, left=Inches(0.7), top=Inches(0.4),
                width=Inches(10), height=Inches(0.6),
                text=title, size=26, bold=True, color=COLOR_INK,
                anchor=MSO_ANCHOR.MIDDLE)
    # 서브타이틀
    if subtitle:
        add_textbox(slide, left=Inches(0.7), top=Inches(0.95),
                    width=Inches(10), height=Inches(0.35),
                    text=subtitle, size=12, color=COLOR_INK_MUTED,
                    anchor=MSO_ANCHOR.TOP)
    # 페이지 번호 (우상단)
    add_textbox(slide, left=Inches(11.2), top=Inches(0.4),
                width=Inches(2), height=Inches(0.4),
                text=f"{page_no} / {total}",
                size=11, color=COLOR_INK_MUTED, align=PP_ALIGN.RIGHT)
    # 푸터 라인
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.05),
                                      Inches(12.83), Inches(7.05))
    line.line.color.rgb = COLOR_BORDER
    line.line.width = Pt(0.5)
    # 푸터 텍스트
    add_textbox(slide, left=Inches(0.5), top=Inches(7.1),
                width=Inches(8), height=Inches(0.3),
                text="지엠티 프로젝트 매니저 — 회사 도입 안내",
                size=9, color=COLOR_INK_MUTED)


def add_bullet(tf, text: str, *, level: int = 0, size: int = 14,
               bold: bool = False, color: RGBColor = COLOR_INK):
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    if not tf.text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = level
    bullet_char = "•" if level == 0 else "–"
    r = p.add_run()
    r.text = f"{bullet_char}  {text}"
    _set_run(r, size=size, bold=bold, color=color)


def add_tip_box(slide, *, left, top, width, height, text: str):
    box = add_filled_rect(slide, left=left, top=top, width=width, height=height,
                          fill=COLOR_TIP_BG, line=COLOR_TIP_BORDER, line_width=1.0)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(150000)
    tf.margin_right = Emu(150000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"💡  {text}"
    _set_run(r, size=12, color=COLOR_INK)


# ────────────────────────────────────────────────────────────────────────────
# 슬라이드 빌더
# ────────────────────────────────────────────────────────────────────────────
TOTAL_PAGES = 10


def build():
    prs = Presentation()
    # 16:9 와이드
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # ── 1. 표지 ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    # 배경
    add_filled_rect(s, left=0, top=0, width=prs.slide_width, height=prs.slide_height,
                    fill=COLOR_BG_LIGHT, line=None)
    # 좌측 강조 띠
    add_filled_rect(s, left=0, top=0, width=Inches(0.5), height=prs.slide_height,
                    fill=COLOR_ACCENT, line=None)
    # 큰 제목
    add_textbox(s, left=Inches(1.2), top=Inches(2.6),
                width=Inches(11), height=Inches(1.2),
                text="지엠티 프로젝트 매니저",
                size=44, bold=True, color=COLOR_INK)
    add_textbox(s, left=Inches(1.2), top=Inches(3.7),
                width=Inches(11), height=Inches(0.6),
                text="사용자 안내",
                size=20, color=COLOR_INK_MUTED)
    # 메타
    add_textbox(s, left=Inches(1.2), top=Inches(6.3),
                width=Inches(11), height=Inches(0.4),
                text="v1.0 · 2026-05-08",
                size=11, color=COLOR_INK_MUTED)

    # ── 2. 화면 구성 ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=2, total=TOTAL_PAGES,
                     title="화면 구성",
                     subtitle="화면은 4개 영역으로 구성됩니다.")
    # 좌측 영역 카드 4개 (번호 라벨로 이미지와 매칭)
    areas = [
        ("①  헤더", "프로젝트 선택, 뷰 전환, 새 작업 추가, 사용자 메뉴."),
        ("②  뷰 탭", "표 / 표+간트 / 간트 / 칸반 / 마인드맵 / 대시보드."),
        ("③  본문", "선택한 뷰의 화면 (작업 표·간트 차트 등)."),
        ("④  사이드바", "현재 화면의 단축키 안내 (우측 [?] 버튼)."),
    ]
    card_top = Inches(1.5)
    for i, (head, body) in enumerate(areas):
        top = card_top + Inches(i * 1.25)
        box = add_filled_rect(s, left=Inches(0.5), top=top,
                              width=Inches(3.8), height=Inches(1.1),
                              fill=COLOR_WHITE, line=COLOR_BORDER, line_width=0.75)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(150000)
        tf.margin_top = Emu(100000)
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        _set_run(r1, size=14, bold=True, color=COLOR_ACCENT_DARK)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        _set_run(r2, size=11, color=COLOR_INK)
    # 우측 이미지 자리
    add_image_placeholder(s, left=Inches(4.7), top=Inches(1.5),
                          width=Inches(8.2), height=Inches(5.4),
                          label="📷  이미지 1 — 전체 화면",
                          hint="브라우저 전체 화면 캡처. ①헤더 ②뷰 탭 ③본문 ④사이드바를 캡처 후 직접 번호로 표시 가능.")

    # ── 3. 가입과 로그인 ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=3, total=TOTAL_PAGES,
                     title="가입과 로그인",
                     subtitle="회사 이메일로 가입 후 관리자 승인 → 로그인")
    # 3단계
    steps = [
        ("①", "회원가입", "회사 이메일·비밀번호 입력 후\n[가입] 클릭."),
        ("②", "관리자 승인", "관리자가 [회원 관리]에서\n승인하면 사용 가능."),
        ("③", "로그인", "승인 후 로그인 시\n조직 회원 정보 자동 연결."),
    ]
    for i, (num, head, body) in enumerate(steps):
        left = Inches(0.5 + i * 4.2)
        # 숫자 배지
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(1.5),
                                    Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        _set_run(cr, size=20, bold=True, color=COLOR_WHITE)
        # 헤딩
        add_textbox(s, left=left + Inches(0.85), top=Inches(1.45),
                    width=Inches(3.2), height=Inches(0.4),
                    text=head, size=16, bold=True, color=COLOR_INK)
        # 본문
        add_textbox(s, left=left + Inches(0.85), top=Inches(1.85),
                    width=Inches(3.2), height=Inches(1.0),
                    text=body, size=11, color=COLOR_INK_MUTED)
    # 이미지
    add_image_placeholder(s, left=Inches(0.5), top=Inches(3.4),
                          width=Inches(12.4), height=Inches(3.3),
                          label="📷  이미지 2 — 로그인·가입 화면",
                          hint="좌: 로그인 화면 (이메일·비밀번호 입력란).  우: 가입 직후 첫 화면 (승인 대기 안내 또는 프로젝트 선택 화면).")

    # ── 4. 프로젝트 만들기 ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=4, total=TOTAL_PAGES,
                     title="프로젝트 만들기",
                     subtitle="헤더 [+ 프로젝트] → 모달에서 입력 → [저장]")
    # 좌측 안내
    add_textbox(s, left=Inches(0.5), top=Inches(1.5),
                width=Inches(5.5), height=Inches(0.4),
                text="입력 항목",
                size=15, bold=True, color=COLOR_ACCENT_DARK)
    bullets_tb = s.shapes.add_textbox(Inches(0.5), Inches(2.0),
                                      Inches(5.8), Inches(2.8))
    tf = bullets_tb.text_frame
    tf.word_wrap = True
    bullets = [
        "프로젝트 이름 (필수)",
        "시작일 / 종료일 — 작업 일정의 허용 범위",
        "공수 단위 — 분 / 시간 / 일 / 주 중 선택",
        "투입인원 — 이름 입력 후 Enter로 다음 인원 추가",
        "주간보고용 정보 — 약어·전체명·구분·주관기관 등 (선택)",
    ]
    for i, t in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"•  {t}"
        _set_run(r, size=13, color=COLOR_INK)
        p.space_after = Pt(8)
    # Tip
    add_tip_box(s, left=Inches(0.5), top=Inches(5.3),
                width=Inches(5.8), height=Inches(1.5),
                text="투입인원 입력란에서 이름을 치면 조직 회원 목록이 자동완성으로 표시됩니다. 일부 글자만 입력해도 매칭되는 이름이 좁혀집니다.")
    # 우측 이미지
    add_image_placeholder(s, left=Inches(6.6), top=Inches(1.5),
                          width=Inches(6.3), height=Inches(5.3),
                          label="📷  이미지 3 — 프로젝트 추가 모달",
                          hint="[+ 프로젝트] 클릭 후 모달이 열린 상태. 필수(이름) + 기본 정보(기간·공수단위) + 투입인원 영역이 보이게.")

    # ── 5. 작업 추가 ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=5, total=TOTAL_PAGES,
                     title="작업 추가",
                     subtitle="표 상단의 청색 입력 행에 작업명 입력 → Enter")
    # 좌측 단축키 안내
    add_textbox(s, left=Inches(0.5), top=Inches(1.5),
                width=Inches(5.5), height=Inches(0.4),
                text="표 작업명 셀에서 사용하는 단축키",
                size=15, bold=True, color=COLOR_ACCENT_DARK)
    rows = [
        ("Enter", "같은 레벨 형제로 추가"),
        ("Shift+Enter", "현재 행 위에 형제 추가"),
        ("Insert", "현재 행의 자식 추가"),
        ("Tab", "들여쓰기 (자식으로 만들기)"),
        ("Shift+Tab", "내어쓰기 (부모의 형제로)"),
    ]
    table_top = Inches(2.0)
    for i, (k, v) in enumerate(rows):
        top = table_top + Inches(i * 0.55)
        # 단축키 박스
        kb = add_filled_rect(s, left=Inches(0.5), top=top,
                             width=Inches(1.6), height=Inches(0.45),
                             fill=COLOR_INK, line=None)
        ktf = kb.text_frame
        ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ktf.margin_left = Emu(80000)
        kp = ktf.paragraphs[0]
        kp.alignment = PP_ALIGN.CENTER
        kr = kp.add_run()
        kr.text = k
        _set_run(kr, size=11, bold=True, color=COLOR_WHITE, font_name="Consolas")
        # 설명
        add_textbox(s, left=Inches(2.3), top=top,
                    width=Inches(4), height=Inches(0.45),
                    text=v, size=12, color=COLOR_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
    # Tip
    add_tip_box(s, left=Inches(0.5), top=Inches(5.3),
                width=Inches(5.8), height=Inches(1.5),
                text="작업명을 비운 상태로 Enter 또는 다른 곳 클릭 시 행이 자동 삭제됩니다.")
    # 우측 이미지
    add_image_placeholder(s, left=Inches(6.6), top=Inches(1.5),
                          width=Inches(6.3), height=Inches(5.3),
                          label="📷  이미지 4 — 표 상단 [+ 새 작업 추가] 행",
                          hint="컬럼 헤더 바로 아래에 항상 표시되는 청색 입력 행 캡처.")

    # ── 6. 작업 편집 ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=6, total=TOTAL_PAGES,
                     title="작업 편집",
                     subtitle="셀을 두 번 클릭하면 입력 모드로 전환됩니다.")
    # 좌측 흐름
    flow_steps = [
        ("1번 클릭", "행 선택 (점선 테두리 표시)", COLOR_INK_MUTED),
        ("2번 클릭", "입력 모드 진입 (F2 또는 Enter도 가능)", COLOR_ACCENT),
        ("Enter", "값 확정. 같은 셀 머무르며 ←/→로 다음 셀 이동", COLOR_INK_MUTED),
        ("Esc", "입력 취소", COLOR_INK_MUTED),
    ]
    for i, (k, v, color) in enumerate(flow_steps):
        top = Inches(1.7 + i * 0.7)
        kb = add_filled_rect(s, left=Inches(0.5), top=top,
                             width=Inches(1.4), height=Inches(0.5),
                             fill=color, line=None)
        ktf = kb.text_frame
        ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        kp = ktf.paragraphs[0]
        kp.alignment = PP_ALIGN.CENTER
        kr = kp.add_run()
        kr.text = k
        _set_run(kr, size=11, bold=True, color=COLOR_WHITE)
        add_textbox(s, left=Inches(2.1), top=top,
                    width=Inches(4.5), height=Inches(0.5),
                    text=v, size=12, color=COLOR_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
    # Tip
    add_tip_box(s, left=Inches(0.5), top=Inches(5.3),
                width=Inches(5.8), height=Inches(1.5),
                text="담당자 셀 클릭 시 조직 회원 자동완성이 표시됩니다. 이름 일부를 입력하면 매칭되는 인원으로 좁혀지며, 부서·직위가 함께 노출됩니다.")
    # 우측 이미지
    add_image_placeholder(s, left=Inches(6.6), top=Inches(1.5),
                          width=Inches(6.3), height=Inches(5.3),
                          label="📷  이미지 5 — 셀 편집 (날짜·담당자)",
                          hint="시작일 셀 입력 모드 (날짜 picker 열림) 또는 담당자 셀 자동완성 드롭다운 캡처.")

    # ── 7. 6가지 화면 ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=7, total=TOTAL_PAGES,
                     title="6가지 화면",
                     subtitle="헤더 가운데 탭으로 화면을 전환합니다. 같은 데이터를 다른 방식으로 표시.")
    views = [
        ("📋", "표만", "작업 목록·속성 편집"),
        ("📊", "표+간트", "표와 간트를 동시 표시"),
        ("📈", "간트만", "일정 막대로 표시"),
        ("📌", "칸반", "상태별 카드로 표시"),
        ("🌳", "마인드맵", "트리 노드로 표시"),
        ("📉", "대시보드", "프로젝트별 통계 표시"),
    ]
    cols = 3
    cell_w = Inches(4.0)
    cell_h = Inches(2.4)
    grid_left = Inches(0.7)
    grid_top = Inches(1.6)
    gap = Inches(0.2)
    for i, (icon, name, desc) in enumerate(views):
        col = i % cols
        row = i // cols
        left = grid_left + col * (cell_w + gap)
        top = grid_top + row * (cell_h + gap)
        box = add_filled_rect(s, left=left, top=top,
                              width=cell_w, height=cell_h,
                              fill=COLOR_WHITE, line=COLOR_BORDER, line_width=1.0)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(200000)
        tf.margin_top = Emu(150000)
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = f"{icon}  {name}"
        _set_run(r1, size=18, bold=True, color=COLOR_ACCENT_DARK)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        _set_run(r2, size=11, color=COLOR_INK_MUTED)
        # 작은 캡처 자리
        add_image_placeholder(
            s,
            left=left + Inches(0.2), top=top + Inches(1.0),
            width=cell_w - Inches(0.4), height=Inches(1.3),
            label=f"📷  이미지 {6 + i}", hint=f"{name} 화면 캡처",
        )

    # ── 8. 공유와 권한 ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=8, total=TOTAL_PAGES,
                     title="공유와 권한 부여",
                     subtitle="헤더 ⋮ → [공유]. 권한 부여는 3가지 방법.")
    # 좌측 안내
    add_textbox(s, left=Inches(0.5), top=Inches(1.5),
                width=Inches(5.5), height=Inches(0.4),
                text="공유 방법",
                size=15, bold=True, color=COLOR_ACCENT_DARK)
    methods = [
        ("초대 링크", "[초대 링크 생성] → 복사 → 전달.\n링크로 접속·로그인 시 편집 권한 자동 부여 (7일 유효)."),
        ("가입 회원 일괄 추가", "조직 필터로 좁히기 → 다중 체크 →\n역할(보기/편집) 선택 → [선택 추가]."),
        ("미가입자 사전 등록", "이름·이메일로 등록.\n해당 인원 가입 시 자동으로 멤버로 추가됨."),
    ]
    for i, (head, body) in enumerate(methods):
        top = Inches(2.1 + i * 1.4)
        # 라벨 배지
        badge = add_filled_rect(s, left=Inches(0.5), top=top,
                                width=Inches(0.4), height=Inches(0.4),
                                fill=COLOR_ACCENT, line=None)
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = str(i + 1)
        _set_run(br, size=12, bold=True, color=COLOR_WHITE)
        # 헤딩
        add_textbox(s, left=Inches(1.0), top=top - Inches(0.05),
                    width=Inches(5), height=Inches(0.4),
                    text=head, size=14, bold=True, color=COLOR_INK)
        # 본문
        add_textbox(s, left=Inches(1.0), top=top + Inches(0.35),
                    width=Inches(5), height=Inches(0.9),
                    text=body, size=11, color=COLOR_INK_MUTED)
    # 우측 이미지
    add_image_placeholder(s, left=Inches(6.6), top=Inches(1.5),
                          width=Inches(6.3), height=Inches(5.3),
                          label="📷  이미지 12 — 공유 모달",
                          hint="공유 모달이 열린 상태. 멤버 목록 + 가입 대기(노란 카드) + 이름·이메일로 권한 부여 영역이 모두 보이게.")

    # ── 9. 단축키 ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=9, total=TOTAL_PAGES,
                     title="단축키",
                     subtitle="표 화면에서 자주 사용하는 단축키 목록")
    keys = [
        ("Ctrl+K", "전체 작업 검색 (다른 프로젝트도)"),
        ("Enter", "같은 레벨 새 작업 추가 (작업명 편집 중)"),
        ("Tab / Shift+Tab", "들여쓰기 / 내어쓰기"),
        ("Insert", "현재 행의 자식 추가"),
        ("F2", "포커스된 셀 편집 시작"),
        ("Alt+↑ / Alt+↓", "형제 작업 순서 위/아래로"),
        ("Shift+→ / ←", "트리 펼치기 / 접기"),
        ("Ctrl+Z / Ctrl+Shift+Z", "되돌리기 / 다시"),
        ("Space (행 포커스)", "체크박스 토글"),
        ("Delete", "선택 작업 삭제"),
    ]
    # 2열 배치
    col_w = Inches(6.0)
    row_h = Inches(0.5)
    start_top = Inches(1.7)
    for i, (k, v) in enumerate(keys):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * (col_w + Inches(0.4))
        top = start_top + row * row_h
        # 단축키 박스
        kb = add_filled_rect(s, left=left, top=top,
                             width=Inches(2.3), height=Inches(0.42),
                             fill=COLOR_INK, line=None)
        ktf = kb.text_frame
        ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ktf.margin_left = Emu(100000)
        kp = ktf.paragraphs[0]
        kp.alignment = PP_ALIGN.CENTER
        kr = kp.add_run()
        kr.text = k
        _set_run(kr, size=10, bold=True, color=COLOR_WHITE, font_name="Consolas")
        # 설명
        add_textbox(s, left=left + Inches(2.5), top=top,
                    width=Inches(3.4), height=row_h,
                    text=v, size=11, color=COLOR_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
    # 하단 Tip
    add_tip_box(s, left=Inches(0.5), top=Inches(6.3),
                width=Inches(12.3), height=Inches(0.6),
                text="전체 단축키 일람은 화면 우측 [?] 사이드바에서 확인할 수 있습니다.")

    # ── 10. 도움말과 문의 ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_slide_header(s, page_no=10, total=TOTAL_PAGES,
                     title="도움말과 문의",
                     subtitle="추가 안내가 필요한 경우 참고하세요.")
    # 카드 3개 (실용 정보)
    ctas = [
        ("화면 안내", "각 화면 첫 진입 시\n안내 토스트가 자동 표시.\n다시 보려면 ⋮ → [튜토리얼]."),
        ("단축키 사이드바", "화면 우측 [?] 버튼.\n현재 화면에서 사용 가능한\n단축키 일람."),
        ("문의·요청", "관리자 또는\n⋮ → [피드백] 메뉴를 통해\n의견을 남길 수 있습니다."),
    ]
    card_w = Inches(4.0)
    card_h = Inches(3.5)
    for i, (head, body) in enumerate(ctas):
        left = Inches(0.5 + i * 4.3)
        top = Inches(2.0)
        box = add_filled_rect(s, left=left, top=top,
                              width=card_w, height=card_h,
                              fill=COLOR_WHITE, line=COLOR_ACCENT, line_width=2.0)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(250000)
        tf.margin_right = Emu(250000)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = head
        _set_run(r1, size=20, bold=True, color=COLOR_ACCENT_DARK)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        _set_run(r2, size=13, color=COLOR_INK)

    out_dir = Path("docs/manual")
    out_dir.mkdir(parents=True, exist_ok=True)
    # 기존 파일(이미지 삽입됨)을 보호하기 위해 v2로 저장.
    # 사용자는 v2 파일의 텍스트를 참고해서 v1에 직접 반영하거나, v2에 이미지를 다시 옮길 수 있다.
    out_path = out_dir / "지엠티스마트시트_도입안내_v2.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    p = build()
    print(f"Generated: {p}")
